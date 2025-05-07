from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from openai import OpenAI
import os
from io import BytesIO
import random
import re

client = OpenAI(base_url="https://api.groq.com/openai/v1", api_key=os.getenv("GROQ_API_KEY"))

COLOR_PALETTES = [
    {
        "primary": RGBColor(41, 105, 176),
        "secondary": RGBColor(133, 193, 233),
        "accent": RGBColor(46, 134, 193),
        "background": RGBColor(240, 248, 255),
        "text": RGBColor(44, 62, 80)
    },
    {
        "primary": RGBColor(40, 180, 99),
        "secondary": RGBColor(171, 235, 198),
        "accent": RGBColor(23, 165, 137),
        "background": RGBColor(245, 255, 250),
        "text": RGBColor(44, 62, 80)
    },
    {
        "primary": RGBColor(142, 68, 173),
        "secondary": RGBColor(187, 143, 206),
        "accent": RGBColor(155, 89, 182),
        "background": RGBColor(248, 240, 255),
        "text": RGBColor(44, 62, 80)
    }
]

def generate_pitch_deck(idea):
    prompt = f"""
Generate a pitch deck for this startup idea:
{idea}
Include the following slides:
1. Title
2. Problem
3. Solution
4. Market Opportunity
5. Business Model
6. Traction (if any)
7. Team
8. Ask
Return each slide with a title and 3-5 bullet points.
"""
    response = client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[{"role": "user", "content": prompt}]
    )
    return response.choices[0].message.content

def create_pptx(slide_content):
    prs = Presentation()
    palette = random.choice(COLOR_PALETTES)
    slide_blocks = [block for block in slide_content.split("\n\n") if block.strip()]
    
    for i, slide_block in enumerate(slide_blocks):
        lines = slide_block.strip().split("\n")
        if not lines:
            continue
        title = lines[0].strip("# ").strip()
        bullet_points = [point.strip("- ").strip() for point in lines[1:] if point.strip()]
        
        if i == 0:
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = palette["primary"]
            title_shape = slide.shapes.title
            title_shape.text = title
            title_text_frame = title_shape.text_frame
            title_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            for paragraph in title_text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(44)
                    run.font.color.rgb = RGBColor(255, 255, 255)
            if len(bullet_points) > 0:
                subtitle = slide.placeholders[1]
                subtitle.text = bullet_points[0]
                subtitle_text_frame = subtitle.text_frame
                subtitle_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                for paragraph in subtitle_text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(24)
                        run.font.color.rgb = RGBColor(240, 240, 240)
        else:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            left = Inches(0)
            top = Inches(0)
            width = Inches(0.5)
            height = Inches(7.5)
            shape = slide.shapes.add_shape(1, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = palette["accent"]
            shape.line.fill.background()
            title_shape = slide.shapes.title
            title_shape.text = title
            title_text_frame = title_shape.text_frame
            for paragraph in title_text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(36)
                    run.font.color.rgb = palette["primary"]
            content = slide.placeholders[1].text_frame
            content.clear()
            for point in bullet_points:
                p = content.add_paragraph()
                p.level = 0
                plain_text = point.replace('\\*', 'ESCAPED_ASTERISK')
                pattern = r'\*(.*?)\*'
                bold_segments = re.findall(pattern, plain_text)
                if bold_segments:
                    parts = re.split(r'\*.*?\*', plain_text)
                    for i in range(len(parts)):
                        if parts[i]:
                            regular_text = parts[i].replace('ESCAPED_ASTERISK', '*')
                            run = p.add_run()
                            run.text = regular_text
                            run.font.size = Pt(24)
                            run.font.color.rgb = palette["text"]
                            run.font.bold = False
                        if i < len(bold_segments):
                            bold_text = bold_segments[i].replace('ESCAPED_ASTERISK', '*')
                            run = p.add_run()
                            run.text = bold_text
                            run.font.size = Pt(24)
                            run.font.color.rgb = palette["text"]
                            run.font.bold = True
                else:
                    p.text = point.replace('\\*', '*')
                    for run in p.runs:
                        run.font.size = Pt(24)
                        run.font.color.rgb = palette["text"]
    
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output
